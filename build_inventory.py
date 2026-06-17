# -*- coding: utf-8 -*-
"""Generate Deloitte-branded Feature/Function Inventory: Excel workbook + single-slide PPTX.
Feature tuple = (text, ai_powered, show_on_slide). Excel shows ALL; slide shows show_on_slide only."""

# -- Deloitte brand palette ---------------------------------------------
GREEN      = "86BC25"   # Deloitte Green (primary)
NEON       = "86EB22"   # accent
BLUE       = "00A3E0"   # accent
DARK       = "282728"   # dark gray bg
BLACK      = "000000"
WHITE      = "FFFFFF"
LIGHTGREEN = "EAF4D7"   # tint for row grouping
GREYTEXT   = "595959"
FONT       = "Open Sans"  # falls back gracefully if not installed

T, F = True, False

# -- Shared data model: 9 L1 pillars -> L2 features (text, ai, on_slide) -
PILLARS = [
    ("1", "AI Requirements Ingestion",
     "Turn documents, transcripts & notes into structured design - automatically, with source evidence preserved.",
     [
        ("Multi-format intake - DOCX, TXT, CSV, MP3, WAV or pasted text", F, T),
        ("Document typing across 9 source types", F, F),
        ("Targeted extraction scope by design area", F, F),
        ("Automatic entity extraction from source material", T, T),
        ("Per-entity confidence scoring with color-coded thresholds", T, T),
        ("AI clarification loop - multi-round, conflict vs. FYI tagging", T, T),
        ("Staged extraction review with per-item reject", F, T),
        ("Promote approved extractions to a Change Packet", F, T),
        ("Inline source-document preview", F, F),
        ("Soft-cancel & restore documents", F, F),
        ("AI usage transparency - tokens, cost & model per run", T, T),
        ("Evidence Source registry with type/status filters", F, T),
        ("Field-level provenance with snippets + confidence", F, F),
        ("ServiceNow platform-layer extraction - data models, forms & business logic from BRDs", T, T),
        ("Non-blocking extraction with live progress", F, T),
     ]),
    ("2", "Agentic Design Studio",
     "Author and review the complete agentic design: use cases, workflows, agents, tools, guardrails, data sources.",
     [
        ("11-tab design review across all design entities & evidence", F, T),
        ("Inline editing everywhere - edits auto-create change packets", F, T),
        ("Full agent design depth - prompt, supervision, trust, model, memory, risks", F, T),
        ("AI prompt drafting from the agent design", T, T),
        ("Entity linking - use cases<->agents, tools<->agents", F, T),
        ("Auto slug cross-linking across the report", F, F),
        ("Full-application quality audit with findings", T, T),
        ("Inline data-quality badges (AI-detected vs. manual)", T, T),
        ("Requirements-by-mode matrix (advisory / HITL / autonomous)", F, T),
        ("Orphan detection for unlinked requirements", F, F),
        ("Print / Save PDF for stakeholders", F, F),
        ("Data-model, form-design & business-logic design entities", F, T),
        ("Reusable pattern library - patterns, standards, decisions", F, T),
     ]),
    ("3", "Governed Change Management",
     "Every change flows through reviewable, auditable Change Packets and locked version baselines.",
     [
        ("Change Packet queue - search + filter by app, source, risk, status", F, T),
        ("Packet detail - risk, conflict, validation, baseline impact", F, F),
        ("Change-item cards with extraction confidence", T, T),
        ("Approve & apply - materializes changes with counts", F, T),
        ("Reject with optional reason", F, F),
        ("Send Back - re-triggers extraction on the linked doc", T, T),
        ("Split a packet into an independent copy", F, T),
        ("Bulk approval with release-type semantics + notes", F, T),
        ("Per-item approve/reject within a Change Packet", F, T),
        ("Post-apply consistency check - flags terminology drift", T, T),
        ("Baseline lifecycle - draft to production, with locking", F, T),
        ("Version history - old->new, triggering packet, approver", F, F),
        ("Baseline compare - field-level diff vs. previous", F, T),
        ("Guardrails & data sources materialize on approval", F, F),
        ("Field-level audit log - diffs, changed-by, originating packet", F, T),
     ]),
    ("4", "Quality, Traceability & Testing",
     "Prove requirement-to-design-to-test coverage and auto-generate tests - catch gaps before build.",
     [
        ("Acceptance Criteria management - inline-editable, status workflow", F, T),
        ("Test Case authoring - scoped, 5 case types", F, T),
        ("AI test generation - balanced cases across scenario types", T, T),
        ("Coverage matrix - requirement x case-type with gap flags", F, T),
        ("AI coverage inference - auto-suggests test->requirement links", T, T),
        ("Manual link manager - per-requirement link/unlink", F, F),
        ("Requirement traceability on criteria & test cases", F, T),
        ("Validation & Exception queue with KPIs", F, T),
        ("Exception resolution workflow - notes, resolve, re-open", F, T),
     ]),
    ("5", "ServiceNow Round-Trip & Sync",
     "Two-way sync with the live platform: reverse-engineer a running ServiceNow app into governed design - and push design deltas back out.",
     [
        ("ServiceNow inbound sync - AI reverse-engineer -> reconcile -> independent review", T, T),
        ("Apply-mode governance - additive-auto+HITL / confidence-gate / review-all", F, T),
        ("Dry-run sync preview - gated plan: confidence, destructive flags, new/changed/drift", F, T),
        ("Non-destructive auto-apply - safe additive applies; rest queued as Change Packets", F, T),
        ("Source-hash change tracking - re-stamp unchanged for faster re-sync", F, F),
        ("Business-logic materiality gate - per-app min-confidence + disallowed types", T, T),
        ("Outbound SN delta export - push design deltas back to ServiceNow", F, T),
     ]),
    ("6", "Cost Intelligence",
     "Project, track, and manage AI and run cost across the full lifecycle.",
     [
        ("AI binding generation - infers Now Assist skills from workflow design", T, T),
        ("App cost parameters - cost/assist, overage, expansion, entitlement", F, T),
        ("Hierarchical cost breakdown - agent -> workflow -> skill binding", F, T),
        ("Cost metrics - per-period & annual cost, baseline, savings, ROI", F, T),
        ("Transparent assumptions - formula + parameters shown", F, F),
        ("Now Assist rate card - catalog of 132 skills with filter + search", F, T),
        ("Rate card editing - assists-per-unit & category per skill", F, F),
     ]),
    ("7", "Build & Delivery Handoff",
     "Export build-ready specs for Claude Code / ServiceNow with one click.",
     [
        ("Section-selectable export - entities, evidence, architecture, planning", F, T),
        ("Version selection - current live design or any locked baseline", F, T),
        ("Optional AI design review - appended summary, gaps & notes", T, T),
        ("Live preview - entity counts, workflow steps, HITL gates", F, T),
        ("Download build spec - Markdown handoff for Claude Code / ServiceNow", F, T),
     ]),
    ("8", "AI Operations & Configuration",
     "Tune the AI engine - models, reasoning depth, house rules - and report on usage.",
     [
        ("Model-per-role - configure Claude model for each AI function", T, T),
        ("Extended thinking controls - toggle, effort level, max tokens", T, T),
        ("Usage & cost dashboard - runs/tokens/cost, by-model breakdown", T, T),
        ("AI Guidance house rules - scoped best practices injected into the prompt", T, T),
        ("Learning loop - acceptance stats + one-click 'Save as best practice'", T, T),
        ("Report Builder - 6 report types, audience targeting, multi-format", F, T),
     ]),
    ("9", "Platform Foundation",
     "Multi-application, multi-user workspace with role-based access, agent trust and a delivery dashboard.",
     [
        ("Applications registry - create & manage apps by client, with stage tracking", F, T),
        ("Application identity - owner, stage, confidence threshold, ingest scope", F, T),
        ("Reuse scope - declare reusable records and visibility", F, F),
        ("Team management - add members with roles", F, T),
        ("Per-application agent enablement", F, T),
        ("Per-agent trust level (1-5 dial), persisted per application", F, T),
        ("Delivery dashboard - KPI cards + recent-changes feed", F, T),
        ("Post-apply findings banner - alerts on unresolved consistency findings", T, T),
        ("Governance watchlists - missing owners + reuse-review panels", F, T),
     ]),
]

TOTAL_FEATS = sum(len(f) for *_, f in PILLARS)
TOTAL_AI    = sum(1 for *_, f in PILLARS for _, ai, _ in f if ai)

# =======================================================================
# 1) EXCEL  (full inventory - every feature)
# =======================================================================
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

wb = Workbook()
ws = wb.active
ws.title = "Feature Inventory"
ws.sheet_view.showGridLines = False

thin = Side(style="thin", color="D9D9D9")
border = Border(left=thin, right=thin, top=thin, bottom=thin)

ws.merge_cells("A1:D1")
c = ws["A1"]; c.value = "ServiceNow Agentic SDLC Workbench"
c.font = Font(name=FONT, size=20, bold=True, color=WHITE)
c.fill = PatternFill("solid", fgColor=DARK)
c.alignment = Alignment(horizontal="left", vertical="center", indent=1)
ws.row_dimensions[1].height = 34

ws.merge_cells("A2:D2")
c = ws["A2"]; c.value = "Feature & Function Inventory  |  Levels 1 + 2"
c.font = Font(name=FONT, size=12, bold=True, color=GREEN)
c.fill = PatternFill("solid", fgColor=DARK)
c.alignment = Alignment(horizontal="left", vertical="center", indent=1)
ws.row_dimensions[2].height = 22

ws.merge_cells("A3:D3")
c = ws["A3"]; c.value = "Together makes progress"
c.font = Font(name="Times New Roman", size=11, italic=True, color=GREEN)
c.fill = PatternFill("solid", fgColor=DARK)
c.alignment = Alignment(horizontal="left", vertical="center", indent=1)
ws.row_dimensions[3].height = 20

hdr = ["L1 #", "L1 Capability Pillar  /  L2 Feature & Function", "AI", "L1 Value Statement"]
ws.append([]); hr = 5
for i, h in enumerate(hdr, start=1):
    cell = ws.cell(row=hr, column=i, value=h)
    cell.font = Font(name=FONT, size=10, bold=True, color=WHITE)
    cell.fill = PatternFill("solid", fgColor=GREEN)
    cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True, indent=1)
    cell.border = border
ws.row_dimensions[hr].height = 24

r = hr + 1
pillar_rows = []
for num, name, value, feats in PILLARS:
    ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=4)
    cell = ws.cell(row=r, column=1, value=f"{num}.  {name}")
    cell.font = Font(name=FONT, size=12, bold=True, color=WHITE)
    cell.fill = PatternFill("solid", fgColor=DARK)
    cell.alignment = Alignment(horizontal="left", vertical="center", indent=1)
    ws.row_dimensions[r].height = 26
    r += 1; first = r
    for feat, ai, _slide in feats:
        ws.cell(row=r, column=1).border = border
        fcell = ws.cell(row=r, column=2, value=feat)
        fcell.font = Font(name=FONT, size=10, color=BLACK)
        fcell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True, indent=1)
        fcell.border = border
        a = ws.cell(row=r, column=3, value=("AI" if ai else ""))
        a.font = Font(name=FONT, size=10, bold=True, color=("4F7A0F" if ai else BLACK))
        a.alignment = Alignment(horizontal="center", vertical="center")
        a.fill = PatternFill("solid", fgColor=(LIGHTGREEN if ai else WHITE))
        a.border = border
        v = ws.cell(row=r, column=4, value=value if r == first else "")
        v.font = Font(name=FONT, size=9, italic=True, color=GREYTEXT)
        v.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True, indent=1)
        v.border = border
        r += 1
    last = r - 1
    if last >= first:
        ws.merge_cells(start_row=first, start_column=4, end_row=last, end_column=4)
    pillar_rows.append((name, first, last))

ws.column_dimensions["A"].width = 6
ws.column_dimensions["B"].width = 82
ws.column_dimensions["C"].width = 6
ws.column_dimensions["D"].width = 40
ws.freeze_panes = "A6"
ws.sheet_properties.tabColor = GREEN

# Sheet 2: L1 overview with live counts
ws2 = wb.create_sheet("L1 Overview")
ws2.sheet_view.showGridLines = False
ws2.merge_cells("A1:D1")
c = ws2["A1"]; c.value = "L1 Capability Pillars - Overview"
c.font = Font(name=FONT, size=16, bold=True, color=WHITE)
c.fill = PatternFill("solid", fgColor=DARK)
c.alignment = Alignment(horizontal="left", vertical="center", indent=1)
ws2.row_dimensions[1].height = 30
for i, h in enumerate(["#", "Capability Pillar", "L2 Features", "AI-Powered"], start=1):
    cell = ws2.cell(row=2, column=i, value=h)
    cell.font = Font(name=FONT, size=10, bold=True, color=WHITE)
    cell.fill = PatternFill("solid", fgColor=GREEN)
    cell.alignment = Alignment(horizontal=("left" if i == 2 else "center"), vertical="center", indent=1)
    cell.border = border
ws2.row_dimensions[2].height = 22
rr = 3
for (num, name, value, feats), (label, first, last) in zip(PILLARS, pillar_rows):
    ws2.cell(row=rr, column=1, value=int(num)).alignment = Alignment(horizontal="center", vertical="center")
    nm = ws2.cell(row=rr, column=2, value=name)
    nm.font = Font(name=FONT, size=10, bold=True, color=BLACK)
    nm.alignment = Alignment(horizontal="left", vertical="center", indent=1)
    ws2.cell(row=rr, column=3, value=f"=COUNTA('Feature Inventory'!B{first}:B{last})")
    ws2.cell(row=rr, column=4, value=f'=COUNTIF(\'Feature Inventory\'!C{first}:C{last},"AI")')
    for col in range(1, 5):
        cell = ws2.cell(row=rr, column=col); cell.border = border
        if col >= 3:
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.font = Font(name=FONT, size=10, color=BLACK)
        cell.fill = PatternFill("solid", fgColor=(WHITE if rr % 2 else "F5F9EE"))
    ws2.cell(row=rr, column=1).font = Font(name=FONT, size=10, color=BLACK)
    rr += 1
tot = ws2.cell(row=rr, column=2, value="Total")
tot.font = Font(name=FONT, size=10, bold=True, color=WHITE)
tot.fill = PatternFill("solid", fgColor=DARK)
tot.alignment = Alignment(horizontal="left", vertical="center", indent=1)
ws2.cell(row=rr, column=1).fill = PatternFill("solid", fgColor=DARK)
for col, formula in [(3, f"=SUM(C3:C{rr-1})"), (4, f"=SUM(D3:D{rr-1})")]:
    cell = ws2.cell(row=rr, column=col, value=formula)
    cell.font = Font(name=FONT, size=10, bold=True, color=WHITE)
    cell.fill = PatternFill("solid", fgColor=DARK)
    cell.alignment = Alignment(horizontal="center", vertical="center")
ws2.column_dimensions["A"].width = 5
ws2.column_dimensions["B"].width = 42
ws2.column_dimensions["C"].width = 14
ws2.column_dimensions["D"].width = 14
ws2.sheet_properties.tabColor = DARK
wb.move_sheet("L1 Overview", -1)

xlsx_path = "SDLC Workbench - Feature Inventory.xlsx"
wb.save(xlsx_path)
print("Saved", xlsx_path, "| full features:", TOTAL_FEATS, "| AI:", TOTAL_AI)

# =======================================================================
# 2) POWERPOINT  (single slide, trimmed to show_on_slide features)
# =======================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def rgb(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])
SW, SH = prs.slide_width, prs.slide_height

bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = rgb(DARK)
bg.line.fill.background(); bg.shadow.inherit = False

def textbox(l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    return tb, tf

tb, tf = textbox(Inches(0.45), Inches(0.18), Inches(11.5), Inches(0.9))
p = tf.paragraphs[0]
r1 = p.add_run(); r1.text = "Agentic SDLC Workbench - Feature "
r1.font.name = FONT; r1.font.size = Pt(26); r1.font.bold = True; r1.font.color.rgb = rgb(WHITE)
r2 = p.add_run(); r2.text = "Inventory"
r2.font.name = "Times New Roman"; r2.font.size = Pt(26); r2.font.italic = True; r2.font.color.rgb = rgb(NEON)
p2 = tf.add_paragraph()
rp = p2.add_run()
rp.text = f"9 capability pillars  -  {TOTAL_FEATS} functions  -  {TOTAL_AI} AI-powered"
rp.font.name = FONT; rp.font.size = Pt(11); rp.font.color.rgb = rgb(GREEN)

tb, tf = textbox(Inches(9.7), Inches(0.28), Inches(3.2), Inches(0.4))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
r = p.add_run(); r.text = "Together makes progress"
r.font.name = "Times New Roman"; r.font.size = Pt(12); r.font.italic = True; r.font.color.rgb = rgb(GREEN)

cols, rows = 3, 3
margin_l, margin_r = Inches(0.4), Inches(0.4)
top = Inches(1.25); bottom_pad = Inches(0.25); gut = Inches(0.18)
grid_w = SW - margin_l - margin_r; grid_h = SH - top - bottom_pad
cw = Emu(int((grid_w - gut*(cols-1)) / cols))
ch = Emu(int((grid_h - gut*(rows-1)) / rows))

for idx, (num, name, value, feats) in enumerate(PILLARS):
    cr, cc = divmod(idx, cols)
    x = Emu(int(margin_l + cc*(cw+gut))); y = Emu(int(top + cr*(ch+gut)))
    card = slide.shapes.add_shape(1, x, y, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb = rgb("33322F")
    card.line.color.rgb = rgb(GREEN); card.line.width = Pt(1); card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(7); tf.margin_right = Pt(6); tf.margin_top = Pt(5); tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    rn = p.add_run(); rn.text = f"{num}  "
    rn.font.name = FONT; rn.font.size = Pt(13); rn.font.bold = True; rn.font.color.rgb = rgb(NEON)
    rt = p.add_run(); rt.text = name
    rt.font.name = FONT; rt.font.size = Pt(11); rt.font.bold = True; rt.font.color.rgb = rgb(WHITE)
    p.space_after = Pt(3)
    shown = [(fe, ai) for fe, ai, sl in feats if sl]
    for feat, ai in shown:
        short = feat.split(" - ")[0].split(" (")[0]
        bp = tf.add_paragraph()
        rb = bp.add_run(); rb.text = ">  "
        rb.font.name = FONT; rb.font.size = Pt(8); rb.font.color.rgb = rgb(GREEN)
        rf = bp.add_run(); rf.text = short
        rf.font.name = FONT; rf.font.size = Pt(8); rf.font.color.rgb = rgb("E8E8E8")
        if ai:
            ra = bp.add_run(); ra.text = "  AI"
            ra.font.name = FONT; ra.font.size = Pt(7); ra.font.bold = True; ra.font.color.rgb = rgb(NEON)
        bp.space_after = Pt(0)

pptx_path = "SDLC Workbench - Feature Inventory.pptx"
prs.save(pptx_path)
slide_counts = {num: sum(1 for *_, sl in f if sl) for num, _, _, f in PILLARS}
print("Saved", pptx_path, "| slide bullets/pillar:", slide_counts)
